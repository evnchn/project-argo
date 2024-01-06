import os
import glob
import collections
import collections.abc
from pptx import Presentation

def makelink(a, bfp, endpoint):
    return f"<a target='scherzo' href='{endpoint+bfp}#page={a}'>'{a}'</a>"

def search_pptx_files(word, output_stream, http_mode=False, endpoint=""):
    current_dir = os.getcwd()
    pptx_files = glob.glob(os.path.join(current_dir, "*.pptx"))

    word = word.lower()  # Convert the word to lowercase for case-insensitive matching

    for file_path in pptx_files:
        presentation = Presentation(file_path)
        total_pages = len(presentation.slides)

        found_pages = []  # List to store the page numbers where the word is found

        for page_number, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if word in shape.text.lower():
                        found_pages.append(page_number)

        found_pages = sorted(list(set(found_pages)))
        if http_mode:
            bfp2 = os.path.basename(file_path).replace(".pptx",".pdf")
            if found_pages:
                output_stream.write(f"Word '{word}' found in <a target='scherzo' href='{endpoint+bfp2}'>'{bfp2}'</a>, on page(s): {','.join(makelink(a, bfp2, endpoint) for a in map(str, found_pages))}<br>")
            else:
                output_stream.write(f"Word '{word}' not found in <a target='scherzo' href='{endpoint+bfp2}'>'{bfp2}'</a>.<br>")
        else:
            if found_pages:
                output_stream.write(f"Word '{word}' found in '{os.path.basename(file_path)}', on page(s): {','.join(map(str, found_pages))}\n")
            else:
                output_stream.write(f"Word '{word}' not found in '{os.path.basename(file_path)}'.\n")