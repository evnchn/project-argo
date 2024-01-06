import_string = \
'''
from pptx import Presentation # python-pptx
'''

### LINK START! (https://github.com/evnchn/linkstart.py)
for line in import_string.splitlines():
    if "import" in line:
        print(line)
        try:
            exec(line)
        except:
            if "#" in line:
                package_name = line.split("#")[-1]
            else:
                splits = line.split("import")
                if "from" in line:
                    package_name = splits[0].replace("from","")
                else:
                    package_name = splits[1]
            package_name = package_name.strip()
            print("Installing {}...".format(package_name))    
            import subprocess
            import sys
            subprocess.check_call([sys.executable, "-m", "pip", "install", package_name])
            try:
                exec(line)
            except:
                print("Failed to install {}".format(package_name))
### DONE

import os
import glob
import collections
import collections.abc
from pptx import Presentation

def search_pptx_files(word):
    current_dir = os.getcwd()
    pptx_files = glob.glob(os.path.join(current_dir, "*.pptx"))

    word = word.lower()  # Convert the word to lowercase for case-insensitive matching

    for file_path in pptx_files:
        presentation = Presentation(file_path)
        total_pages = len(presentation.slides)

        found_pages = []  # List to store the page numbers where the word is found

        for page_number, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if word in shape.text.lower():
                        found_pages.append(page_number)

        if found_pages:
            print(f"Word '{word}' found in '{os.path.basename(file_path)}', on page(s): {','.join(map(str, found_pages))}")
        else:
            print(f"Word '{word}' not found in '{os.path.basename(file_path)}'.")

search_word = input("Enter the word to search: ")
search_pptx_files(search_word)